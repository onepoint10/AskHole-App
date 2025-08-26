"""
File Conversion Utility Module
Handles conversion of Office documents to PDF format for AI processing
"""

import os
import sys
import tempfile
import subprocess
import shutil
from pathlib import Path
import logging
from typing import Optional, Tuple

# Office document processing libraries
try:
    from docx import Document
    from openpyxl import load_workbook
    from pptx import Presentation
    OFFICE_LIBS_AVAILABLE = True
except ImportError:
    OFFICE_LIBS_AVAILABLE = False
    logging.warning("Office document libraries not available. Document conversion will be limited.")

# PDF generation (ReportLab)
try:
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image as RLImage
    from reportlab.lib import colors
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False
    logging.warning("ReportLab not available. Falling back to LibreOffice or HTML.")

# LibreOffice conversion (fallback method)
try:
    subprocess.run(['libreoffice', '--version'], capture_output=True, check=True)
    LIBREOFFICE_AVAILABLE = True
except (subprocess.CalledProcessError, FileNotFoundError):
    LIBREOFFICE_AVAILABLE = False
    logging.warning("LibreOffice not available. Using Python libraries for conversion.")

# Optional: MS Word COM on Windows for best fidelity
MSWORD_AVAILABLE = False
if sys.platform.startswith('win'):
    try:
        import pythoncom  # type: ignore
        import win32com.client  # type: ignore
        MSWORD_AVAILABLE = True
    except Exception:
        MSWORD_AVAILABLE = False
        logging.warning("pywin32 not available. To enable high-fidelity DOCX->PDF on Windows, pip install pywin32.")


def _register_cyrillic_font_if_possible():
    """Register a font that supports Cyrillic so ReportLab renders text correctly."""
    if not REPORTLAB_AVAILABLE:
        return None
    candidates = []
    if sys.platform.startswith('win'):
        candidates.extend([
            r"C:\\Windows\\Fonts\\arial.ttf",
            r"C:\\Windows\\Fonts\\ARIALUNI.TTF",  # Arial Unicode MS
            r"C:\\Windows\\Fonts\\segoeui.ttf",
        ])
    else:
        candidates.extend([
            "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
            "/usr/share/fonts/truetype/noto/NotoSans-Regular.ttf",
            "/usr/share/fonts/truetype/freefont/FreeSans.ttf",
        ])
    for path in candidates:
        try:
            if os.path.exists(path):
                pdfmetrics.registerFont(TTFont('DocBodyFont', path))
                return 'DocBodyFont'
        except Exception as e:
            logging.warning(f"Failed to register font {path}: {e}")
    # Fallback to built-in Helvetica (may not render Cyrillic correctly)
    return None


class FileConverter:
    """Handles conversion of various file formats to PDF"""
    
    @staticmethod
    def convert_to_pdf(file_path: str, output_dir: str = None) -> Optional[str]:
        if not os.path.exists(file_path):
            logging.error(f"File not found: {file_path}")
            return None
        
        file_ext = Path(file_path).suffix.lower()
        if file_ext == '.pdf':
            return file_path
        
        # Prefer MS Word on Windows for highest fidelity
        if file_ext in ['.doc', '.docx'] and MSWORD_AVAILABLE:
            pdf = FileConverter._convert_with_msword(file_path, output_dir)
            if pdf:
                return pdf
        
        # Prefer dedicated python conversion + reportlab for Office and text
        if file_ext in ['.doc', '.docx']:
            pdf = FileConverter._convert_doc_docx_to_pdf_python(file_path, output_dir)
            if pdf:
                return pdf
            # fallback
            return FileConverter._convert_with_libreoffice(file_path, output_dir)
        
        if file_ext in ['.xls', '.xlsx']:
            pdf = FileConverter._convert_xls_xlsx_to_pdf_python(file_path, output_dir)
            if pdf:
                return pdf
            return FileConverter._convert_with_libreoffice(file_path, output_dir)
        
        if file_ext in ['.ppt', '.pptx']:
            pdf = FileConverter._convert_ppt_pptx_to_pdf_python(file_path, output_dir)
            if pdf:
                return pdf
            return FileConverter._convert_with_libreoffice(file_path, output_dir)
        
        if file_ext in ['.txt', '.md', '.py', '.js', '.html', '.css', '.xml', '.json', '.csv']:
            return FileConverter._convert_text_to_pdf(file_path, output_dir)
        
        if file_ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp']:
            return FileConverter._convert_image_to_pdf(file_path, output_dir)
        
        logging.warning(f"Unsupported file type for conversion: {file_ext}")
        return None

    @staticmethod
    def _output_pdf_path(file_path: str, output_dir: Optional[str]) -> str:
        if output_dir is None:
            output_dir = os.path.dirname(file_path)
        base = os.path.splitext(os.path.basename(file_path))[0]
        return os.path.join(output_dir, f"{base}.pdf")

    @staticmethod
    def _convert_with_msword(file_path: str, output_dir: Optional[str]) -> Optional[str]:
        if not MSWORD_AVAILABLE:
            return None
        try:
            pdf_path = FileConverter._output_pdf_path(file_path, output_dir)
            pythoncom.CoInitialize()
            word = win32com.client.Dispatch('Word.Application')
            word.Visible = False
            doc = word.Documents.Open(file_path)
            wdFormatPDF = 17
            doc.ExportAsFixedFormat(OutputFileName=pdf_path, ExportFormat=wdFormatPDF)
            doc.Close(False)
            word.Quit()
            pythoncom.CoUninitialize()
            return pdf_path if os.path.exists(pdf_path) else None
        except Exception as e:
            logging.error(f"MS Word COM conversion failed for {file_path}: {e}")
            try:
                word.Quit()
            except Exception:
                pass
            return None

    @staticmethod
    def _convert_doc_docx_to_pdf_python(file_path: str, output_dir: Optional[str]) -> Optional[str]:
        if not (OFFICE_LIBS_AVAILABLE and REPORTLAB_AVAILABLE):
            return None
        try:
            pdf_path = FileConverter._output_pdf_path(file_path, output_dir)
            doc = Document(file_path)

            # Register font for Cyrillic if possible
            body_font = _register_cyrillic_font_if_possible()

            pdf_doc = SimpleDocTemplate(pdf_path, pagesize=letter)
            styles = getSampleStyleSheet()
            if body_font:
                styles['Normal'].fontName = body_font
                styles['Heading1'].fontName = body_font
                styles['Heading2'].fontName = body_font
            story = []

            title_style = ParagraphStyle('ConvertedTitle', parent=styles['Heading1'], fontSize=16, spaceAfter=12)
            story.append(Paragraph(f"Converted from: {os.path.basename(file_path)}", title_style))
            story.append(Spacer(1, 8))

            # Render body paragraphs
            for paragraph in doc.paragraphs:
                text = paragraph.text
                if text and text.strip():
                    story.append(Paragraph(text, styles['Normal']))
                    story.append(Spacer(1, 4))

            # Render tables (basic text-only rendering)
            if doc.tables:
                story.append(Spacer(1, 8))
                story.append(Paragraph("Tables:", styles['Heading2']))
                for t in doc.tables:
                    data = []
                    for row in t.rows:
                        cells = []
                        for cell in row.cells:
                            cell_text = '\n'.join(p.text for p in cell.paragraphs)
                            cells.append(cell_text or '')
                        data.append(cells)
                    tbl = Table(data, repeatRows=1)
                    tbl.setStyle(TableStyle([
                        ('GRID', (0,0), (-1,-1), 0.25, colors.grey),
                        ('VALIGN', (0,0), (-1,-1), 'TOP'),
                    ]))
                    story.append(tbl)
                    story.append(Spacer(1, 6))

            # Attempt to include images referenced in the document (best effort)
            try:
                related = getattr(doc.part, 'related_parts', {})
                if related:
                    story.append(Spacer(1, 8))
                    story.append(Paragraph("Images:", styles['Heading2']))
                    for rel_id, part in related.items():
                        if hasattr(part, 'content_type') and 'image' in part.content_type:
                            img_bytes = part.blob
                            tmp = tempfile.NamedTemporaryFile(delete=False, suffix=Path(part.partname).suffix)
                            tmp.write(img_bytes)
                            tmp.close()
                            try:
                                story.append(RLImage(tmp.name, width=400, preserveAspectRatio=True, mask='auto'))
                                story.append(Spacer(1, 6))
                            finally:
                                # Defer deletion until after build to keep file available
                                pass
            except Exception as e:
                logging.warning(f"Image extraction warning: {e}")

            pdf_doc.build(story)
            return pdf_path if os.path.exists(pdf_path) else None
        except Exception as e:
            logging.error(f"DOC/DOCX to PDF (python) error for {file_path}: {e}")
            return None

    @staticmethod
    def _convert_xls_xlsx_to_pdf_python(file_path: str, output_dir: Optional[str]) -> Optional[str]:
        if not (OFFICE_LIBS_AVAILABLE and REPORTLAB_AVAILABLE):
            return None
        try:
            pdf_path = FileConverter._output_pdf_path(file_path, output_dir)
            wb = load_workbook(file_path, read_only=True)

            body_font = _register_cyrillic_font_if_possible()
            pdf_doc = SimpleDocTemplate(pdf_path, pagesize=letter)
            styles = getSampleStyleSheet()
            if body_font:
                styles['Normal'].fontName = body_font
                styles['Heading1'].fontName = body_font
                styles['Heading2'].fontName = body_font
            story = []

            title_style = ParagraphStyle('ConvertedTitle', parent=styles['Heading1'], fontSize=16, spaceAfter=12)
            story.append(Paragraph(f"Converted from: {os.path.basename(file_path)}", title_style))
            story.append(Spacer(1, 8))

            for sheet_name in wb.sheetnames:
                story.append(Paragraph(f"Sheet: {sheet_name}", styles['Heading2']))
                sheet = wb[sheet_name]
                data = []
                for row in sheet.iter_rows(values_only=True):
                    data.append(["" if v is None else str(v) for v in row])
                if data:
                    tbl = Table(data, repeatRows=1)
                    tbl.setStyle(TableStyle([
                        ('GRID', (0,0), (-1,-1), 0.25, colors.grey),
                        ('VALIGN', (0,0), (-1,-1), 'TOP'),
                    ]))
                    story.append(tbl)
                story.append(Spacer(1, 6))
            wb.close()

            pdf_doc.build(story)
            return pdf_path if os.path.exists(pdf_path) else None
        except Exception as e:
            logging.error(f"XLS/XLSX to PDF (python) error for {file_path}: {e}")
            return None

    @staticmethod
    def _convert_ppt_pptx_to_pdf_python(file_path: str, output_dir: Optional[str]) -> Optional[str]:
        if not (OFFICE_LIBS_AVAILABLE and REPORTLAB_AVAILABLE):
            return None
        try:
            pdf_path = FileConverter._output_pdf_path(file_path, output_dir)
            prs = Presentation(file_path)

            body_font = _register_cyrillic_font_if_possible()
            pdf_doc = SimpleDocTemplate(pdf_path, pagesize=letter)
            styles = getSampleStyleSheet()
            if body_font:
                styles['Normal'].fontName = body_font
                styles['Heading1'].fontName = body_font
                styles['Heading2'].fontName = body_font
            story = []

            title_style = ParagraphStyle('ConvertedTitle', parent=styles['Heading1'], fontSize=16, spaceAfter=12)
            story.append(Paragraph(f"Converted from: {os.path.basename(file_path)}", title_style))
            story.append(Spacer(1, 8))

            slide_num = 1
            for slide in prs.slides:
                story.append(Paragraph(f"Slide {slide_num}", styles['Heading2']))
                slide_num += 1
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text = (shape.text or '').strip()
                        if text:
                            story.append(Paragraph(text, styles['Normal']))
                story.append(Spacer(1, 8))

            pdf_doc.build(story)
            return pdf_path if os.path.exists(pdf_path) else None
        except Exception as e:
            logging.error(f"PPT/PPTX to PDF (python) error for {file_path}: {e}")
            return None

    @staticmethod
    def _convert_with_libreoffice(file_path: str, output_dir: str = None) -> Optional[str]:
        if not LIBREOFFICE_AVAILABLE:
            return None
        try:
            if output_dir is None:
                output_dir = os.path.dirname(file_path)
            with tempfile.TemporaryDirectory() as temp_dir:
                temp_file = os.path.join(temp_dir, os.path.basename(file_path))
                shutil.copy2(file_path, temp_file)
                cmd = ['libreoffice', '--headless', '--convert-to', 'pdf', '--outdir', temp_dir, temp_file]
                result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
                if result.returncode == 0:
                    pdf_name = os.path.splitext(os.path.basename(file_path))[0] + '.pdf'
                    temp_pdf = os.path.join(temp_dir, pdf_name)
                    if os.path.exists(temp_pdf):
                        output_pdf = os.path.join(output_dir, pdf_name)
                        shutil.move(temp_pdf, output_pdf)
                        logging.info(f"Successfully converted {file_path} to {output_pdf} via LibreOffice")
                        return output_pdf
                logging.error(f"LibreOffice conversion failed (code {result.returncode}): {result.stderr}")
                return None
        except subprocess.TimeoutExpired:
            logging.error(f"LibreOffice conversion timed out for {file_path}")
            return None
        except Exception as e:
            logging.error(f"LibreOffice conversion error for {file_path}: {e}")
            return None
    
    @staticmethod
    def _convert_text_to_pdf(file_path: str, output_dir: str = None) -> Optional[str]:
        """Convert text files to PDF"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            return FileConverter._convert_text_to_pdf_from_content(content, file_path, output_dir)
        except Exception as e:
            logging.error(f"Text file conversion error for {file_path}: {e}")
            return None
    
    @staticmethod
    def _convert_text_to_pdf_from_content(content: str, original_file: str, output_dir: str = None) -> Optional[str]:
        """Convert text content to PDF"""
        try:
            if output_dir is None:
                output_dir = os.path.dirname(original_file)
            
            # Create a simple HTML representation
            html_content = f"""
            <!DOCTYPE html>
            <html>
            <head>
                <meta charset="utf-8">
                <title>{os.path.basename(original_file)}</title>
                <style>
                    body {{ font-family: Arial, sans-serif; margin: 40px; line-height: 1.6; }}
                    pre {{ background-color: #f5f5f5; padding: 15px; border-radius: 5px; overflow-x: auto; }}
                    .filename {{ color: #666; font-size: 12px; margin-bottom: 20px; }}
                </style>
            </head>
            <body>
                <div class="filename">File: {os.path.basename(original_file)}</div>
                <pre>{content}</pre>
            </body>
            </html>
            """
            
            # Save HTML temporarily
            html_path = os.path.join(output_dir, f"{os.path.splitext(os.path.basename(original_file))[0]}.html")
            with open(html_path, 'w', encoding='utf-8') as f:
                f.write(html_content)
            
            # Convert HTML to PDF using wkhtmltopdf if available
            try:
                pdf_path = os.path.join(output_dir, f"{os.path.splitext(os.path.basename(original_file))[0]}.pdf")
                
                # Try wkhtmltopdf
                cmd = ['wkhtmltopdf', '--quiet', html_path, pdf_path]
                result = subprocess.run(cmd, capture_output=True, timeout=30)
                
                if result.returncode == 0 and os.path.exists(pdf_path):
                    # Clean up HTML file
                    os.remove(html_path)
                    return pdf_path
                    
            except (subprocess.TimeoutExpired, FileNotFoundError):
                pass
            
            # Fallback: return HTML file (can be processed as text)
            logging.warning(f"PDF conversion not available, returning HTML: {html_path}")
            return html_path
            
        except Exception as e:
            logging.error(f"Text to PDF conversion error: {e}")
            return None
    
    @staticmethod
    def _convert_image_to_pdf(file_path: str, output_dir: str = None) -> Optional[str]:
        """Convert images to PDF"""
        try:
            from PIL import Image
            
            if output_dir is None:
                output_dir = os.path.dirname(file_path)
            
            # Open image
            img = Image.open(file_path)
            
            # Convert to RGB if necessary
            if img.mode != 'RGB':
                img = img.convert('RGB')
            
            # Create PDF path
            pdf_path = os.path.join(output_dir, f"{os.path.splitext(os.path.basename(file_path))[0]}.pdf")
            
            # Save as PDF
            img.save(pdf_path, 'PDF', resolution=100.0)
            
            return pdf_path
            
        except Exception as e:
            logging.error(f"Image to PDF conversion error for {file_path}: {e}")
            return None
    
    @staticmethod
    def get_supported_formats() -> list:
        """Get list of supported input formats"""
        formats = ['.txt', '.md', '.py', '.js', '.html', '.css', '.xml', '.json', '.csv']
        
        if OFFICE_LIBS_AVAILABLE or LIBREOFFICE_AVAILABLE:
            formats.extend(['.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx'])
        
        formats.extend(['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp'])
        
        return formats